"""
KnowMicro - Document Processing Service
Parse documents, extract text, chunk, and prepare for vector storage.
"""
import io
from pathlib import Path
from typing import List, Dict, Any
from app.config import settings
from app.services.chunking_service import chunking_service


class DocumentService:
    """Parse and process uploaded documents."""

    MAX_SIZE = settings.max_document_size_mb * 1024 * 1024

    async def parse(self, content: bytes, filename: str) -> Dict[str, Any]:
        """
        Parse a document and return extracted text plus metadata.
        Returns: {text, file_type, metadata}
        """
        ext = Path(filename).suffix.lower()

        if len(content) > self.MAX_SIZE:
            raise ValueError(f"文件过大，上限 {self.MAX_SIZE // 1024 // 1024} MB")

        if ext == ".pdf":
            text, meta = self._parse_pdf(content)
        elif ext in (".txt", ".md", ".markdown"):
            text = content.decode("utf-8", errors="replace")
            meta = {"encoding": "utf-8"}
        elif ext == ".docx":
            text, meta = self._parse_docx(content)
        elif ext == ".pptx":
            text, meta = self._parse_pptx(content)
        elif ext == ".ppt":
            raise ValueError("旧版 .ppt 格式暂不支持，请转换为 .pptx 后上传")
        else:
            raise ValueError(f"不支持的文件格式: {ext}")

        if not text or not text.strip():
            raise ValueError("文档内容为空或无法提取文本")

        return {
            "text": text.strip(),
            "file_type": ext.lstrip("."),
            "metadata": meta,
        }

    def _parse_pdf(self, content: bytes) -> tuple:
        """Extract text from PDF."""
        from PyPDF2 import PdfReader
        reader = PdfReader(io.BytesIO(content))
        meta = reader.metadata or {}
        text_parts = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
        return "\n\n".join(text_parts), {
            "pages": len(reader.pages),
            "title": str(meta.get("/Title", "")),
        }

    def _parse_docx(self, content: bytes) -> tuple:
        """Extract text from DOCX."""
        from docx import Document
        doc = Document(io.BytesIO(content))
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        return "\n\n".join(text_parts), {
            "paragraphs": len(doc.paragraphs),
        }

    def _parse_pptx(self, content: bytes) -> tuple:
        """Extract text from PPTX, preserving slide structure."""
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        slides: List[Dict[str, Any]] = []
        all_text: List[str] = []

        for i, slide in enumerate(prs.slides):
            slide_texts: List[str] = []
            title = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            # First text on the slide becomes the title
                            if not title:
                                title = t
                            slide_texts.append(t)
            slide_full = "\n".join(slide_texts)
            slides.append({
                "index": i,
                "title": title or f"幻灯片 {i + 1}",
                "text": slide_full,
                "char_count": len(slide_full),
            })
            if slide_full:
                all_text.append(f"[幻灯片 {i + 1}]\n{slide_full}")

        return "\n\n".join(all_text), {
            "slide_count": len(prs.slides),
            "slides": slides,
        }

    def chunk(
        self, text: str, file_type: str = "txt"
    ) -> List[Dict[str, Any]]:
        """
        Chunk text and return list of {text, index, metadata}.
        """
        chunks = chunking_service.auto_chunk(text, file_type)
        result = []
        for i, chunk_text in enumerate(chunks):
            result.append({
                "index": i,
                "text": chunk_text,
                "char_count": len(chunk_text),
            })
        return result


document_service = DocumentService()
